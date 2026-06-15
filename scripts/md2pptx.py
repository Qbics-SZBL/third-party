#!/usr/bin/env python3
###############################################################################
# Copyright (c) Intel Corporation - All rights reserved.                      #
#                                                                             #
# For information on the license, see the LICENSE file.                       #
# SPDX-License-Identifier: BSD-3-Clause                                       #
###############################################################################
"""Convert Markdown to PowerPoint.

Supports two splitting modes:
  --split=rule   Split on horizontal rules (---).
  --split=heading Split on ## headings.

Auto-detection (default): if the input contains a \\n---\\n
separator, rule mode is used; otherwise heading mode is assumed.

In heading mode, slides are further broken up when they exceed a
content-line budget (overflow splitting).  Continuation slides
receive Roman-numeral suffixes (II, III, ...).  This behaviour is
suppressed in rule mode unless --breakup is given explicitly.
"""
import argparse
import os
import re
import shutil
import subprocess

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

SLIDE_WIDTH = Inches(40 / 3)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

CODE_FONT = "Courier New"
CODE_SIZE = Pt(20)
TABLE_SIZE = Pt(12)
BODY_SIZE = Pt(18)

MARGIN = Inches(0.5)
TITLE_TOP = Inches(0.5)
TITLE_BOX = Inches(1.25)
BODY_TOP = TITLE_TOP + TITLE_BOX + Inches(0.1)
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - MARGIN
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
TABLE_ROW_HEIGHT = Inches(0.5)
CODE_LINE_HEIGHT = Inches(0.35)
TEXT_LINE_HEIGHT = Inches(0.4)
BLOCK_SPACING = Inches(0.6)

MAX_CONTENT_LINES = 15


# ---------------------------------------------------------------------------
# Inline markdown
# ---------------------------------------------------------------------------
def _typographic(text):
    text = text.replace("---", "—")
    text = text.replace("--", "–")
    text = re.sub(r"\\([\\`*_{}[\]()#+\-.!|~])", r"\1", text)
    return text


def parse_inline(text):
    """Parse **bold**, *italic*, and `code` into [(text, style), ...]."""
    result = []
    pos = 0
    for m in re.finditer(r"\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`", text):
        if m.start() > pos:
            result.append((_typographic(text[pos : m.start()]), "normal"))
        if m.group(1) is not None:
            result.append((_typographic(m.group(1)), "bold"))
        elif m.group(2) is not None:
            result.append((_typographic(m.group(2)), "italic"))
        else:
            result.append((m.group(3), "code"))
        pos = m.end()
    if pos < len(text):
        result.append((_typographic(text[pos:]), "normal"))
    return result or [("", "normal")]


def add_runs(para, text, font_size=None, code=False):
    if code:
        run = para.add_run()
        run.text = text
        run.font.name = CODE_FONT
        if font_size:
            run.font.size = font_size
        return
    for content, style in parse_inline(text):
        run = para.add_run()
        run.text = content
        if font_size:
            run.font.size = font_size
        if style == "bold":
            run.font.bold = True
        elif style == "italic":
            run.font.italic = True
        elif style == "code":
            run.font.name = CODE_FONT
            if font_size:
                run.font.size = font_size


def no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    pPr.set("indent", "0")
    pPr.set("marL", "0")
    for tag in ("buNone", "buChar", "buAutoNum"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buNone"))


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------
def parse_markdown(path, split, breakup):
    with open(path) as f:
        raw = f.read()
    if split == "auto":
        split = "rule" if "\n---\n" in raw else "heading"
    if split == "rule":
        parts = re.split(r"\n---\n", raw)
    else:
        parts = _split_on_headings(raw)
    slides = [s for s in (_parse_slide(p.strip()) for p in parts) if s]
    slides = [s for s in slides if s["is_title"] or s["blocks"]]
    slides = _split_title_body(slides)
    if breakup or (breakup is None and split == "heading"):
        slides = _breakup_slides(slides)
    return slides


def _split_title_body(slides):
    """Separate title slides that carry heavy body content into title + content."""
    result = []
    for s in slides:
        if not s["is_title"] or not s["blocks"]:
            result.append(s)
            continue
        heavy = any(
            b["type"] in ("code", "table", "bullets", "image")
            or (b["type"] == "text" and len(b["lines"]) > 2)
            for b in s["blocks"]
        )
        if not heavy:
            result.append(s)
            continue
        result.append(
            {
                "title": s["title"],
                "subtitle": s["subtitle"],
                "is_title": True,
                "blocks": [],
                "notes": s.get("notes", ""),
            }
        )
        result.append(
            {
                "title": s["title"],
                "subtitle": "",
                "is_title": False,
                "blocks": s["blocks"],
                "notes": "",
            }
        )
    return result


def _split_on_headings(raw):
    """Split on #, ##, or ### headings, keeping each heading with its body."""
    parts = []
    cur = []
    in_code = False
    for line in raw.split("\n"):
        if line.strip().startswith("```"):
            in_code = not in_code
        if not in_code and re.match(r"^#{1,3} (?!#)", line) and cur:
            parts.append("\n".join(cur))
            cur = []
        cur.append(line)
    if cur:
        parts.append("\n".join(cur))
    return parts


def _block_lines(block):
    """Estimate the number of content lines a block occupies."""
    btype = block["type"]
    if btype == "code":
        return len(block["lines"]) + 1
    if btype == "bullets":
        return len(block["items"])
    if btype == "table":
        return len(block["rows"]) + 1
    if btype == "text":
        nchars = sum(len(l) for l in block["lines"]) + len(block["lines"]) - 1
        return max(1, (nchars + 79) // 80)
    if btype == "heading":
        return 2
    if btype == "image":
        return 8
    return 1


_ROMAN = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X"]


def _breakup_slides(slides):
    """Split slides that exceed MAX_CONTENT_LINES at block boundaries."""
    result = []
    for slide in slides:
        total = sum(_block_lines(b) for b in slide["blocks"])
        if total <= MAX_CONTENT_LINES or slide["is_title"]:
            result.append(slide)
            continue
        chunks = []
        cur_blocks = []
        cur_lines = 0
        for block in slide["blocks"]:
            bl = _block_lines(block)
            if cur_blocks and cur_lines + bl > MAX_CONTENT_LINES:
                chunks.append(cur_blocks)
                cur_blocks = []
                cur_lines = 0
            cur_blocks.append(block)
            cur_lines += bl
        if cur_blocks:
            chunks.append(cur_blocks)
        if len(chunks) <= 1:
            result.append(slide)
            continue
        base_title = slide["title"]
        for idx, chunk in enumerate(chunks):
            numeral = _ROMAN[idx] if idx < len(_ROMAN) else str(idx + 1)
            title = base_title if idx == 0 else f"{base_title} ({numeral})"
            result.append(
                {
                    "title": title,
                    "subtitle": slide["subtitle"] if idx == 0 else "",
                    "is_title": False,
                    "blocks": chunk,
                    "notes": slide.get("notes", "") if idx == 0 else "",
                }
            )
    return result


def _skip_blank(lines, i):
    while i < len(lines) and not lines[i].strip():
        i += 1
    return i


def _parse_slide(text):
    if not text:
        return None
    lines = text.split("\n")
    slide = {
        "title": "", "subtitle": "", "is_title": False, "blocks": [],
        "notes": "",
    }
    i = _skip_blank(lines, 0)

    h1 = False
    if i < len(lines) and re.match(r"^# (?!#)", lines[i]):
        slide["title"] = lines[i][2:].strip()
        h1 = True
        i = _skip_blank(lines, i + 1)
        if i < len(lines) and lines[i].startswith("## "):
            slide["subtitle"] = lines[i][3:].strip()
            i += 1
    elif i < len(lines) and lines[i].startswith("## "):
        slide["title"] = lines[i][3:].strip()
        i += 1
    elif i < len(lines) and lines[i].startswith("### "):
        slide["title"] = lines[i][4:].strip()
        i += 1

    cur = None
    in_code = False
    note_lines = []
    in_note = False
    while i < len(lines):
        line = lines[i]

        if in_note:
            if not line.strip():
                in_note = False
            else:
                note_lines.append(line.strip())
            i += 1
            continue

        m_note = re.match(r"^Note:\s*(.*)$", line)
        if m_note and not in_code:
            if cur:
                slide["blocks"].append(cur)
                cur = None
            in_note = True
            if m_note.group(1):
                note_lines.append(m_note.group(1))
            i += 1
            continue

        if line.strip().startswith("```"):
            if in_code:
                slide["blocks"].append(cur)
                cur, in_code = None, False
            else:
                if cur:
                    slide["blocks"].append(cur)
                    cur = None
                cur = {"type": "code", "lines": []}
                in_code = True
            i += 1
            continue

        if in_code:
            cur["lines"].append(line)
            i += 1
            continue

        if line.strip().startswith("|") and "|" in line.strip()[1:]:
            if cur and cur["type"] != "table":
                slide["blocks"].append(cur)
                cur = None
            if not cur:
                cur = {"type": "table", "rows": []}
            if not re.match(r"^\s*\|[\s\-:|]+\|\s*$", line):
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                cur["rows"].append(cells)
            i += 1
            continue

        if line.startswith("### "):
            if cur:
                slide["blocks"].append(cur)
                cur = None
            slide["blocks"].append({"type": "heading", "text": line[4:].strip()})
            i += 1
            continue

        m_img = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", line)
        if m_img:
            if cur:
                slide["blocks"].append(cur)
                cur = None
            slide["blocks"].append(
                {"type": "image", "alt": m_img.group(1), "path": m_img.group(2)}
            )
            i += 1
            continue

        m_bullet = re.match(r"^\s*(?:\\?[-*])\s+(.*)$", line)
        m_enum = re.match(r"^\s*(\d+)[.)]\s+(.*)$", line)
        if m_bullet or m_enum:
            if cur and cur["type"] != "bullets":
                slide["blocks"].append(cur)
                cur = None
            if not cur:
                cur = {"type": "bullets", "items": []}
            if m_enum:
                cur["items"].append(f"{m_enum.group(1)}. {m_enum.group(2)}")
            else:
                cur["items"].append(m_bullet.group(1))
            i += 1
            continue

        if not line.strip():
            if cur:
                slide["blocks"].append(cur)
                cur = None
            i += 1
            continue

        if cur and cur["type"] == "text":
            cur["lines"].append(line.strip())
        else:
            if cur:
                slide["blocks"].append(cur)
            cur = {"type": "text", "lines": [line.strip()]}
        i += 1

    if cur:
        slide["blocks"].append(cur)
    if h1:
        slide["is_title"] = True
    if note_lines:
        slide["notes"] = " ".join(note_lines)
    return slide


# ---------------------------------------------------------------------------
# Presentation building
# ---------------------------------------------------------------------------
def _find_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return None


def _set_notes(slide, text):
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _image_size(img_path, max_width, max_height):
    """Compute scaled dimensions preserving aspect ratio."""
    if PILImage is not None:
        with PILImage.open(img_path) as img:
            w_px, h_px = img.size
    else:
        w_px, h_px = 1280, 720
    aspect = w_px / h_px
    width = max_width
    height = int(width / aspect)
    if height > max_height:
        height = max_height
        width = int(height * aspect)
    return width, height


def _apply_dark_theme(prs):
    """Set black background on all slide layouts and master."""
    for master in prs.slide_masters:
        bg = master.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
    for layout in prs.slide_layouts:
        bg = layout.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0, 0, 0)


def _set_text_color(prs):
    """Set all text runs to white across every slide."""
    white = RGBColor(255, 255, 255)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = white
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = white


def build_presentation(slides, base_dir="."):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    _apply_dark_theme(prs)
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx <= 1:
                ph.left = MARGIN
                ph.width = CONTENT_WIDTH
        if layout.name != "Title Slide":
            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0:
                    ph.top = TITLE_TOP
                    ph.height = TITLE_BOX
                elif idx == 1:
                    ph.top = BODY_TOP
                    ph.height = BODY_HEIGHT
    for s in slides:
        if s["is_title"]:
            _title_slide(prs, s)
        elif any(b["type"] == "table" for b in s["blocks"]):
            _table_slide(prs, s, base_dir)
        else:
            _content_slide(prs, s, base_dir)
    _set_text_color(prs)
    return prs


def _title_slide(prs, data):
    layout = _find_layout(prs, "Title Slide") or prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = data["title"]
    tf = slide.placeholders[1].text_frame
    if data["subtitle"]:
        tf.text = data["subtitle"]
    for block in data["blocks"]:
        if block["type"] == "text":
            p = tf.add_paragraph()
            add_runs(p, " ".join(block["lines"]))
    _set_notes(slide, data.get("notes", ""))


def _content_slide(prs, data, base_dir="."):
    has_image = any(b["type"] == "image" for b in data["blocks"])
    layout = (
        _find_layout(prs, "Title Only") or prs.slide_layouts[5]
        if has_image
        else _find_layout(prs, "Title and Content") or prs.slide_layouts[1]
    )
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = data["title"]
    if has_image:
        top = BODY_TOP
        for block in data["blocks"]:
            if block["type"] == "image":
                top = _add_image(slide, block, base_dir, MARGIN, top)
            else:
                n = 1
                if block["type"] in ("code", "bullets"):
                    n = len(block.get("lines", block.get("items", [])))
                elif block["type"] == "text":
                    n = len(block["lines"])
                lh = (
                    CODE_LINE_HEIGHT if block["type"] == "code"
                    else TEXT_LINE_HEIGHT
                )
                height = lh * max(n, 1)
                txbox = slide.shapes.add_textbox(MARGIN, top, CONTENT_WIDTH, height)
                tf = txbox.text_frame
                tf.word_wrap = True
                _render_blocks(tf, [block], placeholder=False)
                top += height + BLOCK_SPACING
    else:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        _render_blocks(tf, data["blocks"], placeholder=True)
    _set_notes(slide, data.get("notes", ""))


def _table_slide(prs, data, base_dir="."):
    layout = _find_layout(prs, "Title Only") or prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = data["title"]
    left = MARGIN
    top = BODY_TOP
    width = CONTENT_WIDTH

    for block in data["blocks"]:
        if block["type"] == "table":
            top = _add_table(slide, block["rows"], left, top, width)
        elif block["type"] == "image":
            top = _add_image(slide, block, base_dir, left, top)
        else:
            n = 1
            if block["type"] in ("code", "bullets"):
                n = len(block.get("lines", block.get("items", [])))
            elif block["type"] == "text":
                n = len(block["lines"])
            lh = CODE_LINE_HEIGHT if block["type"] == "code" else TEXT_LINE_HEIGHT
            height = lh * max(n, 1)
            txbox = slide.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True
            _render_blocks(tf, [block], placeholder=False)
            top += height + BLOCK_SPACING
    _set_notes(slide, data.get("notes", ""))


def _render_blocks(tf, blocks, placeholder=False):
    first = True
    for block in blocks:
        btype = block["type"]

        if btype == "text":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            add_runs(p, " ".join(block["lines"]))
            p.line_spacing = 1.0
            if placeholder:
                no_bullet(p)

        elif btype == "heading":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = block["text"]
            run.font.bold = True
            if placeholder:
                no_bullet(p)

        elif btype == "bullets":
            for item in block["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                if not placeholder:
                    item = "• " + item
                add_runs(p, item)
                p.line_spacing = 1.0

        elif btype == "code":
            for line in block["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                add_runs(p, line, font_size=CODE_SIZE, code=True)
                if placeholder:
                    no_bullet(p)
                p.space_before = Pt(0)
                p.space_after = Pt(0)


def _add_table(slide, rows, left, top, width):
    if not rows:
        return top
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, TABLE_ROW_HEIGHT * n_rows
    )
    tbl = shape.table
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            if c < n_cols:
                p = tbl.cell(r, c).text_frame.paragraphs[0]
                add_runs(p, cell_text, font_size=TABLE_SIZE)
                if r == 0:
                    for run in p.runs:
                        run.font.bold = True
    return top + TABLE_ROW_HEIGHT * n_rows + BLOCK_SPACING


def _add_image(slide, block, base_dir, left, top):
    img_path = block["path"]
    if not os.path.isabs(img_path):
        img_path = os.path.join(base_dir, img_path)
    if not os.path.isfile(img_path):
        return top
    max_w = CONTENT_WIDTH
    max_h = SLIDE_HEIGHT - top - MARGIN
    width, height = _image_size(img_path, max_w, max_h)
    cx = left + (max_w - width) // 2
    slide.shapes.add_picture(img_path, cx, top, width, height)
    return top + height + BLOCK_SPACING


# ---------------------------------------------------------------------------
# Post-processing
# ---------------------------------------------------------------------------
def autofit(pptx_path):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
                if body_pr is None:
                    continue
                if (
                    body_pr.find(qn("a:normAutofit")) is not None
                    or body_pr.find(qn("a:spAutoFit")) is not None
                    or body_pr.find(qn("a:noAutofit")) is not None
                ):
                    continue
                etree.SubElement(body_pr, qn("a:normAutofit"))
    prs.save(pptx_path)


def resave(pptx_path):
    """Reflow text via PowerPoint COM (WSL/Windows only)."""
    if not shutil.which("powershell.exe"):
        return
    try:
        winpath = subprocess.check_output(
            ["wslpath", "-w", os.path.abspath(pptx_path)], text=True
        ).strip()
    except (FileNotFoundError, subprocess.CalledProcessError):
        return
    try:
        subprocess.check_call(
            [
                "powershell.exe",
                "-NoProfile",
                "-Command",
                (
                    f"$pp = New-Object -ComObject PowerPoint.Application;"
                    f'$pres = $pp.Presentations.Open("{winpath}");'
                    f"foreach ($slide in $pres.Slides) {{"
                    f"  foreach ($shape in $slide.Shapes) {{"
                    f"    if ($shape.HasTextFrame) {{"
                    f"      $shape.TextFrame2.AutoSize = 0;"
                    f"      $shape.TextFrame2.AutoSize = 2"
                    f"    }}"
                    f"  }}"
                    f"}}"
                    f"$pres.Save();"
                    f"$pres.Close();"
                    f"$pp.Quit();"
                    f"[System.Runtime.InteropServices.Marshal]"
                    f"::ReleaseComObject($pp) | Out-Null"
                ),
            ]
        )
    except (FileNotFoundError, subprocess.CalledProcessError):
        pass


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description="Markdown to PowerPoint")
    ap.add_argument("input", help="input .md file")
    ap.add_argument("output", nargs="?", default=None, help="output .pptx")
    ap.add_argument(
        "--split",
        choices=["auto", "rule", "heading"],
        default="auto",
        help="slide splitting mode (default: auto)",
    )
    ap.add_argument(
        "--breakup",
        action="store_true",
        default=None,
        help="break up oversized slides (auto in heading mode)",
    )
    ap.add_argument(
        "--no-breakup",
        action="store_true",
        help="disable overflow splitting even in heading mode",
    )
    ap.add_argument(
        "--no-resave",
        action="store_true",
        help="skip PowerPoint COM reflow step",
    )
    args = ap.parse_args()

    breakup = None
    if args.breakup:
        breakup = True
    elif args.no_breakup:
        breakup = False

    output = args.output
    if output is None:
        output = os.path.splitext(args.input)[0] + ".pptx"

    slides = parse_markdown(args.input, args.split, breakup)
    base_dir = os.path.dirname(os.path.abspath(args.input))
    prs = build_presentation(slides, base_dir)
    prs.save(output)
    autofit(output)
    if not args.no_resave:
        resave(output)


if __name__ == "__main__":
    main()
